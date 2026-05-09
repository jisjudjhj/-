from datetime import date
from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUTPUT_DIR / "电商推荐系统-比赛演示版-24页.pptx"
PTF_PATH = OUTPUT_DIR / "电商推荐系统-比赛演示版-24页.ptf"


def rgb(hex_code):
    s = hex_code.strip().replace("#", "")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


PALETTE = {
    "bg1": rgb("#F8FAFC"),
    "bg2": rgb("#EEF2FF"),
    "bg3": rgb("#ECFEFF"),
    "bg4": rgb("#FFF7ED"),
    "title": rgb("#0F172A"),
    "text": rgb("#334155"),
    "muted": rgb("#64748B"),
    "white": rgb("#FFFFFF"),
    "line": rgb("#CBD5E1"),
    "brand": rgb("#2563EB"),
    "green": rgb("#059669"),
    "orange": rgb("#EA580C"),
    "purple": rgb("#7C3AED"),
    "dark": rgb("#0B1220"),
}


def set_p(paragraph, text, size=18, bold=False, color=None, align=None, font_name="Microsoft YaHei"):
    paragraph.text = text
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or PALETTE["text"]
    if align is not None:
        paragraph.alignment = align


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, idx, total):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.16), Inches(13.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = PALETTE["line"]
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.4), Inches(7.18), Inches(1.8), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    set_p(p, f"{idx}/{total}", size=10, color=PALETTE["muted"], align=PP_ALIGN.RIGHT)


def add_header(slide, title, subtitle="", accent="brand"):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.32), Inches(0.18), Inches(0.78))
    tag.fill.solid()
    tag.fill.fore_color.rgb = PALETTE[accent]
    tag.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.75), Inches(0.28), Inches(11.9), Inches(0.6))
    set_p(tbox.text_frame.paragraphs[0], title, size=30, bold=True, color=PALETTE["title"])
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.78), Inches(0.88), Inches(12.0), Inches(0.3))
        set_p(sbox.text_frame.paragraphs[0], subtitle, size=13, color=PALETTE["muted"])


def add_bullet_card(slide, x, y, w, h, title, bullets, tcolor="brand"):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PALETTE["white"]
    box.line.color.rgb = PALETTE["line"]
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_p(p, title, size=17, bold=True, color=PALETTE[tcolor])
    for item in bullets:
        pp = tf.add_paragraph()
        set_p(pp, f"• {item}", size=14, color=PALETTE["text"])


def add_code_card(slide, x, y, w, h, title, lines):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PALETTE["dark"]
    box.line.color.rgb = rgb("#1E293B")
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_p(p, title, size=12, bold=True, color=rgb("#93C5FD"), font_name="Consolas")
    for line in lines:
        pp = tf.add_paragraph()
        set_p(pp, line, size=11, color=rgb("#E2E8F0"), font_name="Consolas")


def add_cards_row(slide, y, cards):
    x0 = 0.8
    w = 2.9
    h = 1.3
    gap = 0.23
    for i, (name, value, c) in enumerate(cards):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x0 + i * (w + gap)), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = PALETTE["white"]
        box.line.color.rgb = PALETTE["line"]
        tf = box.text_frame
        tf.clear()
        set_p(tf.paragraphs[0], name, size=12, color=PALETTE["muted"])
        pp = tf.add_paragraph()
        set_p(pp, value, size=24, bold=True, color=PALETTE[c])


def add_flow(slide, y, nodes):
    x = 0.8
    w = 2.3
    h = 1.15
    gap = 0.2
    for i, node in enumerate(nodes):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + i * (w + gap)), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = PALETTE["white"]
        box.line.color.rgb = PALETTE["line"]
        tf = box.text_frame
        tf.clear()
        set_p(tf.paragraphs[0], node[0], size=13, bold=True, color=PALETTE["brand"], align=PP_ALIGN.CENTER)
        pp = tf.add_paragraph()
        set_p(pp, node[1], size=11, color=PALETTE["text"], align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            ar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + i * (w + gap) + w - 0.02), Inches(y + 0.44), Inches(0.18), Inches(0.24))
            ar.fill.solid()
            ar.fill.fore_color.rgb = PALETTE["muted"]
            ar.line.fill.background()


def add_table(slide, x, y, w, h, headers, rows):
    shp = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tb = shp.table
    for i, htxt in enumerate(headers):
        c = tb.cell(0, i)
        c.text = htxt
        c.fill.solid()
        c.fill.fore_color.rgb = rgb("#E2E8F0")
        set_p(c.text_frame.paragraphs[0], htxt, size=12, bold=True, color=PALETTE["title"], align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        for cidx, value in enumerate(row):
            c = tb.cell(r, cidx)
            c.text = value
            set_p(c.text_frame.paragraphs[0], value, size=11, color=PALETTE["text"])


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"])
    top = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5))
    top.fill.solid(); top.fill.fore_color.rgb = rgb("#DBEAFE"); top.line.fill.background()
    mid = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.8), Inches(3.8))
    mid.fill.solid(); mid.fill.fore_color.rgb = PALETTE["white"]; mid.line.color.rgb = PALETTE["line"]
    set_p(s.shapes.add_textbox(Inches(1.0), Inches(0.55), Inches(11.5), Inches(0.5)).text_frame.paragraphs[0], "比赛演示 · 技术实现完整版", size=14, bold=True, color=PALETTE["brand"])
    tf = s.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.2)).text_frame
    set_p(tf.paragraphs[0], "大数据电商推荐系统", size=52, bold=True, color=PALETTE["title"])
    p = tf.add_paragraph(); set_p(p, "22页答辩版｜技术架构 + 核心实现 + 系统级优化", size=22, color=PALETTE["brand"])
    p = tf.add_paragraph(); set_p(p, f"日期：{date.today().strftime('%Y-%m-%d')}  |  演示人：________", size=14, color=PALETTE["muted"])

    # 2 Agenda
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "目录", "建议总时长 10~15 分钟", "purple")
    ag = [("01", "项目背景与目标"), ("02", "系统架构与数据链路"), ("03", "算法设计与实现"), ("04", "系统级治理优化"), ("05", "效果评估与演示脚本"), ("06", "创新点与规划")]
    for i, (no, txt) in enumerate(ag):
        add_bullet_card(s, Inches(0.8), Inches(1.5 + i * 0.88), Inches(11.8), Inches(0.72), f"{no}  {txt}", [], "brand")

    # 3-4
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "1. 项目背景与业务痛点", "为什么这个题目有比赛价值")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(4.9), "业务痛点", ["新用户冷启动，推荐易退化热门", "单一算法各有短板", "需要兼顾效果/多样性/稳定性", "比赛要求可演示可解释可评估"], "orange")
    add_bullet_card(s, Inches(6.9), Inches(1.5), Inches(5.7), Inches(4.9), "项目目标", ["搭建在线+离线+实时闭环", "提升CTR/加购率/下单转化", "降低同质化，提升多样性", "形成可复现的工程方案"], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg3"]); add_header(s, "2. 比赛评审维度映射", "把技术点映射到可得分点", "green")
    add_cards_row(s, 1.6, [("效果", "CTR/CVR", "brand"), ("工程", "低时延", "green"), ("创新", "自愈机制", "purple"), ("展示", "可解释", "orange")])
    add_bullet_card(s, Inches(0.8), Inches(3.3), Inches(12.0), Inches(3.1), "评委关注点回答策略", ["不仅给模型，还给系统闭环", "不仅讲原理，还展示接口和代码路径", "不仅讲效果，还展示异常治理与回退"], "brand")

    # 5-8
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "3. 业务闭环架构", "曝光 -> 行为 -> 画像 -> 推荐")
    add_flow(s, 2.1, [("用户触达", "首页/搜索/详情"), ("推荐服务", "Hybrid在线计算"), ("行为回流", "点击/加购/下单"), ("数据沉淀", "行为与事件表"), ("策略优化", "权重与重排")])
    add_bullet_card(s, Inches(0.8), Inches(4.1), Inches(12.0), Inches(2.2), "闭环价值", ["推荐不是一次性结果，而是持续优化系统", "每次曝光都会反哺下一轮推荐策略"], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "4. 技术架构总览", "Spring Boot + MySQL + Redis + MQ + Kafka/Flink")
    add_bullet_card(s, Inches(0.8), Inches(1.45), Inches(12.0), Inches(4.95), "分层结构", ["展示层：管理端推荐预览 + 用户端推荐页面", "服务层：RecommendationServiceImpl / HybridEngine / CF / CB", "数据层：MySQL行为与画像、Redis缓存与实时特征", "消息层：RabbitMQ异步消费、Kafka/Flink流式聚合", "分析层：离线画像/分群/A-B评估"], "purple")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg4"]); add_header(s, "5. 核心数据模型", "可解释与可评估依赖结构化数据", "orange")
    add_table(s, Inches(0.8), Inches(1.6), Inches(12.0), Inches(4.9), ["表名", "作用", "关键字段"], [
        ("user_behavior", "用户行为日志", "user_id, product_id, behavior_type, create_time"),
        ("user_preference", "用户画像", "category_preferences(JSON), tag_preferences(JSON)"),
        ("recommendation_event", "推荐归因", "scene, token, exposure/click/order"),
        ("analytics_recommendation_result", "推荐快照", "snapshot_date, rank_no, score"),
        ("stream_user_category_preference", "实时偏好", "user_id, category_id, preference_score"),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg3"]); add_header(s, "6. 离线与实时双链路", "稳定基线 + 新鲜信号", "green")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(4.9), "离线链路", ["用户相似度矩阵", "商品共现矩阵", "用户画像与分群", "推荐快照生成", "A/B报表汇总"], "brand")
    add_bullet_card(s, Inches(6.9), Inches(1.5), Inches(5.7), Inches(4.9), "实时链路", ["CDC进入Kafka", "Flink聚合用户偏好与热度", "Redis写入实时特征", "在线推荐优先读取实时信号", "异常时自动回退快照"], "orange")

    # 9-13 algorithms & implementation
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "7. 算法总览", "CF + CB + 热门 + Hybrid")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(3.9), Inches(4.9), "User-CF", ["相似用户行为迁移", "适合高活跃用户", "稀疏时需兜底"], "brand")
    add_bullet_card(s, Inches(4.8), Inches(1.5), Inches(3.9), Inches(4.9), "Content-CB", ["标签/品类/价格匹配", "冷启动更稳定", "需防过窄"], "purple")
    add_bullet_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(4.9), "Hybrid", ["按实验组动态权重", "融合排序+探索", "支持解释与A/B"], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "8. CF实现细节（代码级）", "CollaborativeFiltering.java")
    add_code_card(s, Inches(0.7), Inches(1.5), Inches(6.1), Inches(3.8), "User-Based CF", [
        "buildUserVector(userId)",
        "purchase=6, favorite=4, cart=2, view=0.2",
        "sim(u,v)=cosine(userVector, candidateVector)",
        "aggregate weighted product score",
        "fallback: behavior -> hot(diversified)",
    ])
    add_bullet_card(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(3.8), "工程优化", ["Redis缓存用户向量", "时间衰减增强新鲜行为", "降级链路日志化", "热门回退跨类目打散"], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "9. CB实现细节（代码级）", "ContentBasedFiltering.java", "purple")
    add_code_card(s, Inches(0.7), Inches(1.5), Inches(6.1), Inches(3.8), "Content Score", [
        "score=0.35*tagSim + 0.25*categoryWeight",
        "    +0.15*priceMatch + 0.15*salesNorm",
        "    +0.10*ratingNorm",
        "优先读实时偏好，回退保存画像",
        "fallback hot -> diversifyByCategory",
    ])
    add_bullet_card(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(3.8), "为什么稳定", ["冷启动可用", "标签/品类/价格多信号结合", "解释信息可直接展示给评委"], "orange")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg3"]); add_header(s, "10. Hybrid融合 + A/B测试", "HybridRecommendationEngine + ABTestFramework", "green")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(4.9), "融合策略", ["按实验组读取CF/CB/HOT权重", "三路候选合并排序", "CF或CB缺失时自动重分配权重", "记录曝光事件形成评估闭环"], "brand")
    add_bullet_card(s, Inches(6.9), Inches(1.5), Inches(5.7), Inches(4.9), "评估策略", ["控制组：热门基线", "实验组：Hybrid/CF强化/CB强化", "指标：CTR、加购率、下单率、退款率"], "purple")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg4"]); add_header(s, "11. 核心代码路径", "评委追问时的快速定位")
    add_code_card(s, Inches(0.8), Inches(1.6), Inches(12.0), Inches(4.9), "Key Classes", [
        "service/impl/RecommendationServiceImpl.java",
        "recommendation/HybridRecommendationEngine.java",
        "recommendation/CollaborativeFiltering.java",
        "recommendation/ContentBasedFiltering.java",
        "recommendation/UserPreferenceBootstrapService.java",
        "controller/RecommendationController.java",
        "controller/AdminController.java",
    ])

    # 14-18 governance and APIs
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "12. 核心API设计", "用户端 + 管理端", "brand")
    add_table(s, Inches(0.8), Inches(1.6), Inches(12.0), Inches(4.9), ["Method", "Path", "Purpose"], [
        ("GET", "/api/recommendations/personal", "个性化推荐"),
        ("GET", "/api/recommendations/guess-you-like", "猜你喜欢"),
        ("GET", "/api/recommendations/hot", "热门推荐"),
        ("POST", "/api/recommendations/behavior", "行为上报"),
        ("GET", "/api/admin/recommend/preview/{userId}", "后台推荐预览"),
        ("GET", "/api/admin/recommend/compare/{userId}", "多算法对比"),
        ("POST", "/api/admin/recommend/profile-bootstrap", "全量画像重建"),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "13. 系统级治理1：用户画像自愈", "从“补丁”升级到“机制”", "orange")
    add_code_card(s, Inches(0.7), Inches(1.5), Inches(6.2), Inches(3.9), "UserPreferenceBootstrapService", [
        "ensureUserPreferenceInitialized(userId, force)",
        "check missing/invalid profile",
        "merge behavior signal + category prior",
        "build category/tag/price profile",
        "persist to user_preference",
    ])
    add_bullet_card(s, Inches(7.1), Inches(1.5), Inches(5.5), Inches(3.9), "落地范围", ["推荐入口统一调用", "管理端支持全量重建", "启动任务自动补齐", "避免空画像用户被热门主导"], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg3"]); add_header(s, "14. 系统级治理2：反同质化重排", "CF/CB/Hybrid兜底统一多类目打散", "purple")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(2.1), "治理思路", ["热门候选按类目分桶轮询", "限制单类目头部占比", "保证列表多样性同时保留头部质量"], "brand")
    add_code_card(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(2.3), "Diversify Logic", [
        "groupByCategory(candidates)",
        "while result not full: poll one item each bucket",
        "deduplicate by productId",
        "return diversified ranked list",
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "15. 可观测性与回滚能力", "线上稳定性的关键")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9), "观测", ["Actuator + Prometheus", "推荐事件全链路日志", "降级原因可追踪", "数据质量巡检任务"], "green")
    add_bullet_card(s, Inches(6.85), Inches(1.5), Inches(5.8), Inches(4.9), "回滚", ["快照兜底保障可用", "A/B可快速回切控制组", "流量护栏抑制异常结果"], "orange")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg4"]); add_header(s, "16. 召回-排序-重排全链路", "不是单次排序，是多阶段优化")
    add_bullet_card(s, Inches(0.8), Inches(1.7), Inches(3.8), Inches(4.7), "召回层", ["CF召回", "CB召回", "热门兜底"], "brand")
    add_bullet_card(s, Inches(4.9), Inches(1.7), Inches(3.8), Inches(4.7), "排序层", ["融合打分", "实验组权重", "探索扰动"], "purple")
    add_bullet_card(s, Inches(9.0), Inches(1.7), Inches(3.8), Inches(4.7), "重排层", ["类目商家护栏", "近重复抑制", "会话去疲劳"], "green")

    # 19-22 presentation and close
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "17. 推荐可解释性设计", "让评委看到“为什么推荐”", "green")
    add_bullet_card(s, Inches(0.8), Inches(1.6), Inches(12.0), Inches(2.8), "解释机制", ["解释来源：CF命中、标签匹配、品类偏好、热门趋势", "前端展示每个商品的命中信号", "从黑盒推荐升级为可解释推荐"], "brand")
    add_code_card(s, Inches(0.8), Inches(4.7), Inches(12.0), Inches(1.6), "Reason Template", [
        "“你近期关注的{category}类商品中，该SKU在相似用户里转化更高，且当前热度上涨。”"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "18. 性能与稳定性（请替换实测）", "建议赛前填入你环境的真实数据", "orange")
    add_cards_row(s, 1.6, [("P95时延", "<180ms", "brand"), ("可用性", "99.9%+", "green"), ("日行为量", "10w+示例", "purple"), ("回退率", "<1%示例", "orange")])
    add_bullet_card(s, Inches(0.8), Inches(3.4), Inches(12.0), Inches(2.8), "建议展示", ["接口压测：/personal /guess-you-like /admin/recommend/preview", "业务指标：CTR、加购率、下单率对照提升", "稳定性：异常场景下的快照回退命中率"], "brand")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg3"]); add_header(s, "19. 现场演示脚本", "按这个顺序讲最稳", "purple")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(4.9), "演示步骤", [
        "1) 打开推荐预览，先展示Hybrid结果",
        "2) 切换CF/CB/热门，强调同一用户结果差异",
        "3) 展示用户画像与分群，解释命中原因",
        "4) 触发行为上报，刷新结果展示实时变化",
        "5) 调用profile-bootstrap接口展示自愈能力",
        "6) 展示A/B对照与关键指标结论",
    ], "green")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg4"]); add_header(s, "20. 高频问答准备", "提前准备，现场不慌")
    add_bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9), "算法问答", ["Q: 为什么不用单一深度模型？", "A: 比赛阶段优先可解释与工程稳定，Hybrid更稳", "Q: 冷启动怎么处理？", "A: 画像自愈 + 多样性兜底 + 行为快速回流"], "brand")
    add_bullet_card(s, Inches(6.85), Inches(1.5), Inches(5.8), Inches(4.9), "工程问答", ["Q: 如何证明有效？", "A: A/B对照 + 曝光到下单全链路指标", "Q: 异常怎么保障？", "A: 快照回退 + 策略回滚 + 护栏机制"], "orange")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg2"]); add_header(s, "21. 创新点总结", "比赛得分重点页", "green")
    add_bullet_card(s, Inches(0.8), Inches(1.6), Inches(12.0), Inches(4.8), "核心创新", [
        "在线/离线/实时融合：兼顾效果与稳定",
        "推荐解释 + A/B闭环：可验证而非拍脑袋",
        "用户画像自愈：系统级解决画像缺失问题",
        "多样性重排护栏：解决同质化真实痛点",
        "结论：不仅能跑通，还能持续优化与扩展",
    ], "purple")

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, PALETTE["bg1"]); add_header(s, "22. Roadmap 与致谢", "Q&A")
    add_bullet_card(s, Inches(0.8), Inches(1.7), Inches(12.0), Inches(3.0), "后续规划", [
        "短期：补齐真实压测页与业务指标自动报表",
        "中期：引入序列建模，增强会话级兴趣捕捉",
        "长期：探索强化学习重排与策略自动调参",
    ], "brand")
    thanks = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(12.0), Inches(1.2))
    thanks.fill.solid(); thanks.fill.fore_color.rgb = rgb("#DBEAFE"); thanks.line.fill.background()
    tf = thanks.text_frame; tf.clear()
    set_p(tf.paragraphs[0], "谢谢各位评委老师！", size=30, bold=True, color=PALETTE["brand"], align=PP_ALIGN.CENTER)
    p = tf.add_paragraph(); set_p(p, "Q & A", size=18, bold=True, color=PALETTE["purple"], align=PP_ALIGN.CENTER)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)

    prs.save(PPTX_PATH)
    copyfile(PPTX_PATH, PTF_PATH)
    print(f"slides={total}")
    print(str(PPTX_PATH))
    print(str(PTF_PATH))


if __name__ == "__main__":
    build()
